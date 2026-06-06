import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Define file paths
OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
PPTX_PATH = os.path.join(OUTPUT_DIR, "Geopolitical_Oil_Research_Presentation.pptx")
CHART_PATH = os.path.join(OUTPUT_DIR, "oil_geopolitical_timeline_chart.png")
LOGO_PATH = os.path.join(OUTPUT_DIR, "iit_roorkee_logo.png")

# Slides content database
SLIDE_DATA = [
    {
        "title": "Executive Summary: Conflict, Volatility, and Corporate Resilience",
        "bullets": [
            "Geopolitical Decoupling: Modern conflicts inject a persistent risk premium that decouples oil prices from underlying physical market supply/demand balances.",
            "Logistics & Sourcing Redirection: Vessel rerouting around Red Sea/Bab al-Mandab raises ton-mile demand and shipping costs, forcing refiners to restructure supply lines.",
            "Strategic Hedging: Energy-intensive sectors must shift from reactive procurement to proactive, multi-layered risk management (feedstock flexibility, layered hedging).",
            "Pragmatic Sourcing (The India Angle): Emerging economies successfully buffered domestic inflation by importing discounted Russian crude under U.S./G7 waiver frameworks."
        ],
        "notes": "Good morning/afternoon. Today we will examine how modern geopolitical conflicts shape the global oil market. The core takeaway is that oil is not just a physical commodity; it is a geopolitical transmission mechanism. We will trace how conflicts translate into price premiums, look at the logistics bottlenecks of chokepoints like the Red Sea, explore the impact on corporate margins, and analyze how countries like India navigated these dynamics pragmatically to secure their energy needs."
    },
    {
        "title": "Core Research Question & Scope of Inquiry",
        "bullets": [
            "Primary Research Question: 'How do modern geopolitical conflicts alter the crude oil economy through supply expectations, transit bottlenecks, and sanctions—and how should exposed firms redesign pricing, sourcing, and hedging strategies?'",
            "Key Areas of Analysis:\n• Geopolitical transmission channels (price caps, chokepoint transit risks)\n• Differentiating physical supply shocks from sentiment-driven risk premiums\n• Refinery economics and corporate finance (aviation/logistics sector margin transmission)",
            "Temporal Scope: 2022 to May 23, 2026 (capturing Russia-Ukraine sanctions, Red Sea crisis, and 2026 Middle East tensions)."
        ],
        "notes": "Our project is framed around a central question: How should corporate decision-makers respond when war risk interacts with uncertain demand? We are looking at a critical window from 2022 to mid-2026. This period is highly unique because it represents a transition from acute physical supply scares (in early 2022) to complex, politically carved-out markets governed by G7 price caps, shadow tanker fleets, and naval chokepoint blockades."
    },
    {
        "title": "Chronology of Volatility: Brent Crude & Conflict Milestones",
        "bullets": [
            "2022 Spike: Russia-Ukraine invasion pushes Brent to a peak of $133.18/bbl in March 2022 due to immediate supply-loss panic.",
            "2023 Consolidation: G7 price caps ($60/bbl) and Russian supply redirection stabilize prices in the $75-$85 range.",
            "Chokepoint Shocks: Late 2023 Houthi Red Sea attacks and April 2024 Iran-Israel direct strikes inject immediate risk premiums.",
            "2026 Volatility: Combat risk involving Iran in early 2026 revives chokepoint fears, driving Brent up despite predicted annual surpluses."
        ],
        "notes": "This slide visually maps major geopolitical events directly to daily spot prices using the chart shown. You can see the initial massive spike in early 2022. Interestingly, look at the G7 price cap implementation date on December 5, 2022—it actually stabilized prices as trade routes adapted. In late 2023 and 2024, naval escalations in the Red Sea created minor price ripples but major shipping rate spikes. By early 2026, prices surged again as Strait of Hormuz risks resurfaced."
    },
    {
        "title": "How Conflict Translates to Price: The 5 Transmission Channels",
        "bullets": [
            "1. Direct Production Threats: Destruction of pipelines, refineries, or wells (e.g., strikes on Russian refineries, Gulf oil facility vulnerabilities).",
            "2. Transit Chokepoint Impairments: Threat to critical sea lanes (Bab al-Mandab, Strait of Hormuz).",
            "3. Logistics & Shipping Economics: Increased shipping distances (ton-miles), higher hull insurance premiums, and freight rate inflation.",
            "4. Sanctions & Compliance Constraints: G7 price caps, financial compliance burdens, and banking/payment restrictions.",
            "5. Risk Premium & Expectations: Speculative trading and options hedging by financial actors, driving up prices before a single physical barrel is lost."
        ],
        "notes": "To build a structured framework, we must understand how conflict transmits to price. It happens through five distinct channels. While direct production threats are rare, the other four channels—chokepoints, shipping costs, sanctions compliance, and paper market speculation—are constantly active in the modern economy. Even if oil flows normally, the cost of transit and the price of 'risk insurance' in the futures market drives the crude benchmark higher."
    },
    {
        "title": "Decoupling Fundamentals: Physical Supply vs. Sentiment Risk",
        "bullets": [
            "Physical Supply Shock:\n• Actual loss of barrels from the market (e.g., pipeline closures, production shutdowns)\n• Forces immediate physical rationing and sharp, long-term price structural changes.",
            "Risk Premium / Sentiment Decoupling:\n• Oil continues to flow, but markets price in the probability of future disruption\n• Example: April 2024 Iran-Israel strikes. No oil was lost, but Brent spiked due to call-option volume and short-covering.",
            "Implication for Firms: Buying oil during risk-premium spikes locks in inflated costs; distinguishing between the two is key for procurement."
        ],
        "notes": "One of our project objectives is to distinguish between actual supply loss and market-imposed risk premiums. For instance, when Brent crude spiked in April 2024, it was driven by fear, not a physical deficit. In fact, global inventories were relatively balanced. For procurement officers, recognizing when a price increase is a speculative 'fear premium' prevents them from buying at the absolute peak of the market."
    },
    {
        "title": "Redrawing the Energy Map: Chokepoints & Shipping Economics",
        "bullets": [
            "The Red Sea Crisis (Nov 2023): Houthi shipping attacks forced tanker fleets to avoid the Suez Canal/Bab al-Mandab.",
            "The Cape of Good Hope Reroute:\n• Adds 10 to 14 days to voyage times between Asia/Middle East and Europe\n• Increases fuel consumption and operating costs per voyage.",
            "The Ton-Mile Effect: Rerouting dramatically increased global 'ton-mile' demand, soaking up available tanker capacity and inflating global freight rates.",
            "Insurance Premium Spikes: War-risk insurance premiums for Red Sea transits surged, adding hundreds of thousands of dollars to single trips."
        ],
        "notes": "Let's zoom into the logistics channel. When Houthi attacks began in late 2023, the physical flow of oil didn't stop, but it had to take a massive detour around the southern tip of Africa. This rerouting is a prime example of shipping economics shaping crude prices. It increased transit times by up to two weeks. This 'ton-mile' expansion meant the same amount of oil required more tankers, causing shipping rates and fuel costs to skyrocket globally."
    },
    {
        "title": "Industry Focus: Crude Volatility Transmission to Aviation",
        "bullets": [
            "The Exposure: Jet fuel (kerosene) is the single largest operating expense for airlines, representing 30% to 40% of total costs.",
            "Imperfect Pass-Through: Airlines cannot adjust ticket prices instantly due to advance bookings and competition. Spot price spikes result in immediate margin contraction.",
            "Working Capital Strain: Sudden oil spikes require airlines to post more collateral for fuel hedges and pay higher cash amounts for immediate fueling, draining liquidity.",
            "Refining Spread Correlation: The 'crack spread' (cost of crude vs. cost of refined jet fuel) often widens during conflicts due to regional refining deficits, compounding the crude price spike."
        ],
        "notes": "To keep our project concrete, we selected Aviation as our focal industry. Jet fuel is the lifeblood of airline finance. When Brent spikes, jet fuel prices rise even faster. Because airlines sell tickets months in advance, they cannot instantly pass these fuel cost increases to consumers. This 'imperfect pass-through' creates a severe lag, causing airline profitability to dip sharply during periods of high geopolitical tension."
    },
    {
        "title": "India's Energy Security: Sourcing Pragmatism Under Sanctions",
        "bullets": [
            "Pragmatic Sourcing: India is the world's 3rd largest oil importer (importing 85% of its crude). Energy security dictates national policy.",
            "Leveraging the Sanction Discount: Following the G7 price cap, India imported massive volumes of Russian Urals crude, which traded at discounts of up to $20-$30/bbl relative to Brent.",
            "Import Share Surges: Russian oil went from <2% of India's import basket in 2021 to over 33% by March 2026 and 38.6% by May 2026.",
            "Domestic Buffer: Sourcing discounted crude protected India's domestic refining margins and cushioned the public from local fuel inflation."
        ],
        "notes": "India represents a fascinating case study in energy pragmatism. As a country that imports 85% of its crude, India chose energy security over political alignment. By purchasing heavily discounted Russian oil under the waiver and price cap rules, Indian refiners successfully protected domestic margins and shielded the Indian economy from the inflation shocks that struck Western nations."
    },
    {
        "title": "Corporate Defense: Strategic Hedging vs. Speculative Exposure",
        "bullets": [
            "The Hedging Paradox: Hedging is a risk-reduction tool, not a profit center. Locking in prices during a peak is an 'expensive mistake.'",
            "Best Practice: Layered Hedging:\n• Airlines/refiners should hedge systematically (e.g., hedging 50% of volume 12 months out, 25% 6 months out) rather than making all-or-nothing bets.",
            "Flexible Feedstock Contracts: Refiners must build operational flexibility to switch between different crude grades depending on regional discounts (e.g., switching Urals/Middle East grades).",
            "Dynamic Surcharges: Implementing automated fuel surcharges in pricing models to reduce the pass-through lag."
        ],
        "notes": "How should corporate finance respond? The answer lies in hedging and procurement discipline. Companies often make the mistake of buying expensive hedges only after a war starts and prices spike. A sensible approach is systematic, layered hedging—locking in prices incrementally over a rolling 12-month period. Additionally, refiners must maintain feedstock flexibility, allowing them to switch crude sources based on localized discount spreads."
    },
    {
        "title": "Government Policy: Strategic Reserves and Tax Interventions",
        "bullets": [
            "Strategic Petroleum Reserve (SPR) Releases:\n• Provides short-term supply liquidity (e.g., US SPR releases in 2022)\n• Risk: Leaves reserves depleted, reducing future national leverage.",
            "Tax and Excise Adjustments:\n• Temporarily lowering fuel excise duties (e.g., India's fuel tax cuts) shields consumers\n• Risk: Directly reduces government fiscal revenue.",
            "Diplomatic Diversification: Securing state-to-state term contracts (e.g., India-Brazil, India-Gulf relations) to ensure supply continuity during embargoes."
        ],
        "notes": "Finally, let's look at public policy. Governments are not helpless during crude spikes. They can release oil from Strategic Petroleum Reserves to calm the market, or cut domestic fuel taxes to protect consumers. However, these tools are temporary and come at a cost—tax cuts stress government budgets, and releasing strategic reserves leaves the country vulnerable to future, more severe supply interruptions."
    },
    {
        "title": "Operational Roadmap: Actionable Corporate Responses",
        "bullets": [
            "Immediate Actions:\n• Review existing fuel hedging contracts and enforce rolling, layered hedging schedules\n• Embed automatic fuel surcharges into downstream supply contracts.",
            "Medium-Term Actions:\n• Audit logistics supply chains to verify exposure to maritime chokepoints (Red Sea, Strait of Hormuz)\n• Increase storage buffer capacity to manage shipping delays.",
            "Long-Term Strategy:\n• Invest in dual-feedstock refining capabilities\n• For logistics/aviation, accelerate transition to energy-efficient assets and alternative fuels (e.g., Sustainable Aviation Fuel)."
        ],
        "notes": "To conclude, companies cannot control geopolitics, but they can control their exposure. Our main recommendation is to transition from reactive compliance to proactive structural resilience. This means checking supply chain chokepoint exposure today, formalizing a layered hedging program, and investing in the operational flexibility to process different grades of crude depending on where the market discounts lie. Thank you, and I am happy to open the floor to questions."
    }
]

def add_header_decoration(slide):
    """Draws a premium thin red header accent line at the top of the slide"""
    # Create a thin line shape: left = 0.5 inches, top = 1.15 inches, width = 12.33 inches, height = 0.02 inches (1.5pt)
    # Using a rectangle shape for exact pixel control
    from pptx.enum.shapes import MSO_SHAPE
    left = Inches(0.5)
    top = Inches(1.15)
    width = Inches(12.33)
    height = Inches(0.02)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(220, 38, 38) # Crimson Red accent (#dc2626)
    line.line.color.rgb = RGBColor(220, 38, 38)

def add_institute_logo(slide):
    """Adds the IIT Roorkee logo to the top right corner of the slide"""
    if os.path.exists(LOGO_PATH):
        # Position: left = 12.13 inches, top = 0.25 inches, width = 0.75 inches, height = 0.75 inches
        slide.shapes.add_picture(LOGO_PATH, Inches(12.13), Inches(0.25), Inches(0.75), Inches(0.75))
    else:
        print(f"Warning: Logo image not found at {LOGO_PATH}")

def main():
    print("Generating Premium Light-Themed PowerPoint presentation with IIT Roorkee branding...")
    
    # Initialize Presentation
    prs = Presentation()
    
    # Set to modern 16:9 widescreen format
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # --- TITLE SLIDE ---
    # We will build a clean, minimalist light-themed Title Slide
    blank_layout = prs.slide_layouts[6] # Blank Slide for total custom layout control
    slide = prs.slides.add_slide(blank_layout)
    
    # Set slide background color to pure off-white (#f8fafc)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252)
    
    # Place a large Title block in the center
    # left = 0.8 inches, top = 2.0 inches, width = 11.73 inches, height = 3.5 inches
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.73), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Geopolitical Conflicts & The Crude Oil Economy"
    p.font.name = "Arial"
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42) # Slate-900 (#0f172a)
    p.space_after = Pt(20)
    
    p2 = tf.add_paragraph()
    p2.text = "Mapping Benchmark Volatility, Trade Rerouting, and Corporate Sourcing Strategies (2022-2026)"
    p2.font.name = "Arial"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(71, 85, 105) # Slate-600 (#475569)
    p2.space_after = Pt(40)

    p3 = tf.add_paragraph()
    p3.text = "A Research Project Submission"
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = RGBColor(220, 38, 38) # Crimson Red (#dc2626)
    
    # Place IIT Roorkee logo on the Title Slide at the top right
    # Size it slightly larger (1.1 inches) for the Title Slide
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(11.4), Inches(0.5), Inches(1.1), Inches(1.1))
        
    slide.notes_slide.notes_text_frame.text = "Welcome to the presentation. This presentation analyzes how international conflicts shape the crude oil benchmark volatility, impacts trade routing, and alters corporate sourcing strategies, focusing on the period from 2022 to mid-2026."

    # --- CONTENT SLIDES ---
    # We will build clean, structured light-themed slides for the content
    for idx, slide_info in enumerate(SLIDE_DATA):
        slide = prs.slides.add_slide(blank_layout)
        
        # Set slide background to clean white
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Add IIT Roorkee logo on the top right
        add_institute_logo(slide)
        
        # Add a premium red accent bar below the header
        add_header_decoration(slide)
        
        # Set Slide Title: left = 0.5 inches, top = 0.35 inches, width = 10.0 inches
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(10.0), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = Inches(0)
        tf_title.margin_top = Inches(0)
        
        p_title = tf_title.paragraphs[0]
        p_title.text = slide_info["title"]
        p_title.font.name = "Arial"
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(15, 23, 42) # Slate-900 (#0f172a)
        
        # Draw Content
        if idx == 2:
            # Special layout for Slide 3 (Chart Slide)
            # Left: text bullets (width 5.5 inches)
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.5), Inches(5.2))
            tf_content = content_box.text_frame
            tf_content.word_wrap = True
            tf_content.margin_left = Inches(0)
            
            for b_idx, bullet in enumerate(slide_info["bullets"]):
                p = tf_content.add_paragraph() if b_idx > 0 else tf_content.paragraphs[0]
                p.text = "• " + bullet
                p.font.name = "Arial"
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(51, 65, 85) # Slate-700
                p.space_after = Pt(12)
                p.line_spacing = 1.15
                
            # Right: Embed timeline chart
            if os.path.exists(CHART_PATH):
                # Left: 6.25 inches, Top: 1.4 inches, Width: 6.58 inches, Height: 5.2 inches
                slide.shapes.add_picture(CHART_PATH, Inches(6.25), Inches(1.4), Inches(6.58), Inches(5.2))
                print("Embedded timeline chart image on Slide 3.")
            else:
                print(f"Warning: Chart image not found at {CHART_PATH}")
                
        else:
            # Standard Slide Layout: single unified text box with premium spacing
            # Left: 0.5 inches, Top: 1.4 inches, Width: 12.33 inches, Height: 5.2 inches
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.33), Inches(5.2))
            tf_content = content_box.text_frame
            tf_content.word_wrap = True
            tf_content.margin_left = Inches(0)
            
            for b_idx, bullet in enumerate(slide_info["bullets"]):
                p = tf_content.add_paragraph() if b_idx > 0 else tf_content.paragraphs[0]
                p.text = "• " + bullet
                p.font.name = "Arial"
                p.font.size = Pt(15)
                p.font.color.rgb = RGBColor(51, 65, 85) # Slate-700
                p.space_after = Pt(14)
                p.line_spacing = 1.2
                
        # Set Speaker Notes
        slide.notes_slide.notes_text_frame.text = slide_info["notes"]
        
    # Save Presentation
    prs.save(PPTX_PATH)
    print(f"Successfully generated presentation file: {PPTX_PATH}")

if __name__ == "__main__":
    main()
